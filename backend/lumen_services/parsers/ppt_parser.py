"""M38.4 — PowerPoint parser.

Each *slide* becomes one text chunk; ``speaker_notes`` (if any) become
a *separate* chunk with ``kind="ppt_notes"`` and the same page_number
so the KB search layer can scope to a specific slide. Images embedded
in a slide are returned via the ``image_assets`` list — the parser
itself does not write to storage; the caller (``document_tasks`` /
``knowledge.py``) turns those into ``ImageAsset`` rows on commit.

The parser is **authoritative** for chunking (see
``preserves_chunks = True``). Each chunk dict carries the new
M38.4 ``DocumentChunk`` fields directly (modality / page_number /
image_caption) so the chunk-creation path can populate the new
columns without another join.

PowerPoint ``MSO_SHAPE_TYPE.PICTURE`` constant is ``13``. Using
``shape.shape_type == MSO_SHAPE_TYPE.PICTURE`` (rather than the
literal ``13``) keeps the magic-number out of the codebase and
survives python-pptx renumbering.
"""
from __future__ import annotations

from typing import Any, Dict, List

from . import BaseParser


class PPTParser(BaseParser):
    """PowerPoint (.pptx / .ppt) parser — one chunk per slide + notes."""

    #: Skip the secondary text-split pass in ``DocumentParser``.
    preserves_chunks: bool = True

    def get_type(self) -> str:
        return "ppt"

    def parse(self, file_path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:  # pragma: no cover - dependency miss
            return self._fallback_parse(file_path, reason=f"python-pptx not installed: {exc}")

        try:
            presentation = Presentation(file_path)
        except Exception as exc:
            return self._fallback_parse(file_path, reason=f"python-pptx: {type(exc).__name__}: {exc}")

        chunk_records: List[Dict[str, Any]] = []
        image_assets: List[Dict[str, Any]] = []
        slide_count = 0

        # Slides are 1-indexed in user-facing strings (slide 1 is the
        # first slide). Spec §3.2 / §1.3 require ``page_number`` for
        # PPT chunks; ``page_number=1`` must mean the first slide, not
        # index 0. ``enumerate(prs.slides, 1)`` gives us that for free.
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            slide_count += 1

            slide_text_lines: List[str] = []
            image_shape_index = 0
            for shape in slide.shapes:
                # Text-bearing shapes — text_frame.text joins all
                # paragraphs with newlines so a multi-paragraph bullet
                # list renders as separate visual lines.
                if shape.has_text_frame:
                    text = (shape.text_frame.text or "").strip()
                    if text:
                        slide_text_lines.append(text)
                    continue

                # Embedded image (PICTURE == 13). Skip chart / SmartArt
                # / placeholder picture frames — spec §2 defers those.
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_shape_index += 1
                    image_record = _extract_image_shape(shape, slide_idx, image_shape_index)
                    if image_record is not None:
                        image_assets.append(image_record)
                    # Don't put image shape's alt text into the slide
                    # chunk — that conflates visual content with the
                    # slide's main text. (Spec: PPT text chunk is for
                    # text only; image chunks live in image_assets.)

            slide_text = "\n".join(slide_text_lines)
            slide_chunk_text = f"Slide {slide_idx}\n{slide_text}".rstrip()
            chunk_records.append(
                {
                    "content": slide_chunk_text,
                    "chunk_index": len(chunk_records),
                    "length": len(slide_chunk_text),
                    "strategy": "ppt_slide",
                    "modality": "text",
                    "sheet_name": None,
                    "page_number": slide_idx,
                    "image_caption": None,
                    "chunk_metadata": {
                        "kind": "ppt_slide",
                        "page_number": slide_idx,
                        "image_count": image_shape_index,
                    },
                }
            )

            # Speaker notes: spec §1.3 says notes become a *separate*
            # chunk so search can scope "find me notes for slide 3".
            # python-pptx's ``slide.notes_slide`` is None if the slide
            # has no notes slide attached — guard with ``has_notes_slide``
            # which is the documented "notes exist?" probe.
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    notes_chunk_text = f"Slide {slide_idx} Notes\n{notes_text}"
                    chunk_records.append(
                        {
                            "content": notes_chunk_text,
                            "chunk_index": len(chunk_records),
                            "length": len(notes_chunk_text),
                            "strategy": "ppt_notes",
                            "modality": "text",
                            "sheet_name": None,
                            "page_number": slide_idx,
                            "image_caption": None,
                            "chunk_metadata": {
                                "kind": "ppt_notes",
                                "page_number": slide_idx,
                            },
                        }
                    )

        # M38.4 (2026-09-01): emit one ``image`` chunk per extracted
        # picture so downstream search can filter ``modality == 'image'``
        # and the multimodal embedder can index these alongside standalone
        # image uploads. ``content`` carries the slide_dedup_key so the
        # caption text reflects the source position — until the worker
        # LLM-generates richer captions in v2 the dedup key ("slide 3
        # image 2") is the only descriptive anchor we have.
        for image_record in image_assets:
            caption_anchor = image_record.get("slide_dedup_key", "").replace("_", " ")
            chunk_records.append({
                "content": caption_anchor,
                "chunk_index": len(chunk_records),
                "length": len(caption_anchor),
                "strategy": "ppt_image",
                "modality": "image",
                "sheet_name": None,
                "page_number": image_record.get("page_number"),
                "image_caption": caption_anchor,
                "chunk_metadata": {
                    "kind": "ppt_image",
                    "page_number": image_record.get("page_number"),
                    "slide_dedup_key": image_record.get("slide_dedup_key"),
                    "shape_index": image_record.get("shape_index"),
                    "mime": image_record.get("mime"),
                },
            })

        joined_text = "\n\n".join(c["content"] for c in chunk_records)
        return {
            "text": joined_text,
            "metadata": {
                "type": self.get_type(),
                "format": "pptx",
                "slide_count": slide_count,
                "image_count": len(image_assets),
                "file_path": file_path,
                "chunk_count": len(chunk_records),
            },
            "chunks": chunk_records,
            "chunk_metadata": [c["chunk_metadata"] for c in chunk_records],
            "image_assets": image_assets,
        }


def _extract_image_shape(shape: Any, page_number: int, shape_index: int) -> Dict[str, Any] | None:
    """Pull bytes + metadata out of an ``MSOSHAPE_TYPE.PICTURE`` shape.

    Returns ``None`` for shapes that the python-pptx API refuses to
    serialise (rare — happens on broken OOXML with a malformed image
    part). The caller skips ``None`` so a single bad image doesn't
    kill the whole PPT.

    The raw bytes go into the ``image_assets`` entry; the caller is
    responsible for handing them to ``storage.put_object`` and then
    inserting an ``ImageAsset`` row with the resulting ``storage_key``.
    We keep the parser side-effect-free — parsers don't talk to the
    storage backend or the DB.
    """
    try:
        image = shape.image
        image_bytes = image.blob
        content_type = image.content_type or "image/png"
    except Exception:
        # ``shape.image`` can raise on shapes that LOOK like pictures
        # but reference a missing image part. Skip and continue.
        return None

    # Slide-level dedup key. The caller can use this to derive a
    # deterministic storage key (e.g. ``<doc>/slide_<n>_img_<m>.<ext>``)
    # so re-uploading the same PPT twice produces stable keys for
    # idempotency.
    ext = (content_type.split("/")[-1] or "png").lower()
    slide_dedup_key = f"slide_{page_number}_img_{shape_index}"

    return {
        "page_number": page_number,
        "shape_index": shape_index,
        "slide_dedup_key": slide_dedup_key,
        "bytes": image_bytes,
        "mime": content_type,
        "extension": ext,
        "size_bytes": len(image_bytes),
    }


__all__ = ["PPTParser"]
