"""M38.4 Step 4 — multimodal parser factory + DocumentParser integration.

Verifies:

- ``BaseParser.preserves_chunks = False`` by default (legacy 6 parsers
  keep the existing secondary text-split behaviour)
- ``ExcelParser`` / ``PPTParser`` / ``ImageParser`` all override
  ``preserves_chunks = True``
- ``DocumentParserFactory.PARSERS`` includes the 3 new types and
  ``TYPE_PATTERNS`` routes xlsx/pptx/png filenames correctly
- ``DocumentParser.parse()`` honours ``preserves_chunks=True`` — when
  the parser already produced authoritative chunks, the secondary
  text-split is skipped (so per-sheet / per-slide / per-image
  boundaries survive)
- ``DocumentParser.supported_formats`` recognises xlsx / pptx / png
  extensions
- ``DocumentParserFactory.get_available_types()`` lists the 3 new
  multimodal types so the frontend dropdown can render them
"""
from __future__ import annotations

import os

import pytest
from openpyxl import Workbook
from PIL import Image as PILImage
from pptx import Presentation

from lumen_services.document_parser import DocumentParser
from lumen_services.parsers import (
    BaseParser,
    DocumentParserFactory,
    ExcelParser,
    GeneralParser,
    ImageParser,
    PPTParser,
)


# ----------------------------------------------------------------------
# Parser-level class attribute checks
# ----------------------------------------------------------------------


def test_base_parser_preserves_chunks_default_is_false():
    """Legacy 6 parsers don't set the flag → default False → secondary
    text-split still applies (no behaviour change for pre-M38.4 docs)."""
    assert BaseParser.preserves_chunks is False
    assert GeneralParser.preserves_chunks is False


def test_multimodal_parsers_set_preserves_chunks_true():
    assert ExcelParser.preserves_chunks is True
    assert PPTParser.preserves_chunks is True
    assert ImageParser.preserves_chunks is True


# ----------------------------------------------------------------------
# Factory wiring
# ----------------------------------------------------------------------


def test_factory_parsers_dict_includes_multimodal_types():
    assert "excel" in DocumentParserFactory.PARSERS
    assert "ppt" in DocumentParserFactory.PARSERS
    assert "image" in DocumentParserFactory.PARSERS
    assert DocumentParserFactory.PARSERS["excel"] is ExcelParser
    assert DocumentParserFactory.PARSERS["ppt"] is PPTParser
    assert DocumentParserFactory.PARSERS["image"] is ImageParser


def test_factory_detect_doc_type_routes_by_extension():
    """Filename patterns route .xlsx/.xls → "excel", .pptx/.ppt → "ppt",
    .png/.jpg/.jpeg/.webp → "image"."""
    for fname, expected in [
        ("sales.xlsx", "excel"),
        ("data.xls", "excel"),
        ("slides.pptx", "ppt"),
        ("deck.ppt", "ppt"),
        ("logo.png", "image"),
        ("photo.jpg", "image"),
        ("image.jpeg", "image"),
        ("pic.webp", "image"),
        ("screenshot.gif", "image"),
        ("diagram.bmp", "image"),
        ("scan.tiff", "image"),
    ]:
        assert DocumentParserFactory.detect_doc_type(fname) == expected, fname


def test_factory_get_parser_returns_correct_class():
    assert isinstance(DocumentParserFactory.get_parser("excel"), ExcelParser)
    assert isinstance(DocumentParserFactory.get_parser("ppt"), PPTParser)
    assert isinstance(DocumentParserFactory.get_parser("image"), ImageParser)


def test_factory_get_available_types_includes_multimodal():
    """The frontend dropdown reads this list — multimodal types MUST be
    present so users can pick them explicitly."""
    types = {entry["type"] for entry in DocumentParserFactory.get_available_types()}
    assert "excel" in types
    assert "ppt" in types
    assert "image" in types
    # Legacy 6 must also still be present (no regression).
    for legacy in ("general", "paper", "qa", "table", "manual", "laws"):
        assert legacy in types


# ----------------------------------------------------------------------
# DocumentParser integration
# ----------------------------------------------------------------------


def test_document_parser_supports_multimodal_formats():
    """``supported_formats`` must include xlsx / xls / pptx / ppt /
    png / jpg / jpeg / webp / gif / bmp / tiff so the format-family
    classification lands on the multimodal parsers, not 'unknown'."""
    dp = DocumentParser()
    for ext in (".xlsx", ".xls", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".webp"):
        assert ext in dp.supported_formats, ext


def test_document_parser_skips_secondary_split_for_excel(tmp_path):
    """End-to-end: upload .xlsx → ``parse()`` returns the parser's
    authoritative chunks (one per sheet), NOT a secondary text-split
    that would merge them."""
    out = tmp_path / "two_sheets.xlsx"
    wb = Workbook()
    wb.remove(wb.active)
    ws1 = wb.create_sheet("First")
    ws1.append(["a"])
    ws1.append(["1"])
    ws2 = wb.create_sheet("Second")
    ws2.append(["b"])
    ws2.append(["2"])
    wb.save(str(out))

    result = DocumentParser().parse(str(out))

    assert result["metadata"]["type"] == "excel"
    # MUST be 2 chunks (one per sheet), not a single merged text-split
    # chunk that the legacy secondary split would have produced.
    assert len(result["chunks"]) == 2
    assert result["chunks"][0]["sheet_name"] == "First"
    assert result["chunks"][1]["sheet_name"] == "Second"


def test_document_parser_skips_secondary_split_for_ppt(tmp_path):
    out = tmp_path / "two_slides.pptx"
    prs = Presentation()
    for title in ("Alpha", "Beta"):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
    prs.save(str(out))

    result = DocumentParser().parse(str(out))

    assert result["metadata"]["type"] == "ppt"
    assert len(result["chunks"]) == 2
    assert result["chunks"][0]["page_number"] == 1
    assert result["chunks"][1]["page_number"] == 2


def test_document_parser_skips_secondary_split_for_image(tmp_path):
    out = tmp_path / "logo.png"
    PILImage.new("RGB", (40, 40), "red").save(str(out))

    result = DocumentParser().parse(str(out))

    assert result["metadata"]["type"] == "image"
    # Image is exactly 1 chunk — a secondary text-split would break
    # this (e.g. turn "logo" into a 1-char chunk "l").
    assert len(result["chunks"]) == 1
    assert result["chunks"][0]["modality"] == "image"
    assert result["chunks"][0]["image_caption"] == "logo"


def test_document_parser_legacy_pdf_still_uses_secondary_split(tmp_path):
    """Regression: pre-M38.4 PDF uploads must keep the secondary
    text-split behaviour. We can't easily write a real PDF in a
    tmpfile without a 3rd-party dep, so we just confirm the legacy
    flag is False on GeneralParser (which handles PDF)."""
    assert GeneralParser.preserves_chunks is False


def test_document_parser_filename_extension_routes_to_multimodal(tmp_path):
    """A file with .xlsx extension auto-routes to ExcelParser even
    without an explicit ``doc_type`` parameter."""
    out = tmp_path / "auto.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["a", "b"])
    wb.save(str(out))

    result = DocumentParser().parse(str(out))
    assert result["metadata"]["type"] == "excel"
    assert "xlsx" == result["metadata"]["format"]
