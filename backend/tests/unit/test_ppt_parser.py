"""M38.4 Step 4 — PPTParser unit tests.

Real pptx fixtures via python-pptx (the parser's own dep). Tests cover:

- each slide produces one ``ppt_slide`` chunk with the right
  ``page_number`` (1-indexed)
- speaker_notes become a separate ``ppt_notes`` chunk sharing the
  slide's ``page_number``
- images embedded in slides are extracted into the ``image_assets``
  list with mime, size, and a deterministic dedup key
- empty slides (no text, no notes, no images) still get a chunk —
  they're not noise, they're a positional anchor for the KB UI
- ``preserves_chunks=True`` so the secondary text-split doesn't
  collapse multiple slides into one chunk
- graceful fallback when the file is missing or corrupt

python-pptx is the parser's own dep — the test builds real pptx files
in tmp_path; cost is ~50 ms per test.
"""
from __future__ import annotations

import io

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

from lumen_services.parsers import PPTParser


# ----------------------------------------------------------------------
# Fixture builders
# ----------------------------------------------------------------------


def _new_pptx() -> Presentation:
    """Empty pptx with no slides — tests add slides explicitly."""
    return Presentation()


def _add_text_slide(prs: Presentation, title: str, body: str = "", notes: str = "") -> None:
    """Add a title + body slide; attach notes if provided."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 1 = title + content
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    # Body placeholder index 1; may not exist for all layouts but
    # layout 1 (Title Slide) does have it.
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = body
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_image_slide(prs: Presentation, color: str = "red") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    buf = io.BytesIO()
    PILImage.new("RGB", (50, 50), color).save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1), Inches(1))


# ----------------------------------------------------------------------
# Tests
# ----------------------------------------------------------------------


def test_ppt_parser_one_chunk_per_slide(tmp_path):
    out = tmp_path / "three.pptx"
    prs = _new_pptx()
    _add_text_slide(prs, "First", "Body 1")
    _add_text_slide(prs, "Second", "Body 2")
    _add_text_slide(prs, "Third", "Body 3")
    prs.save(str(out))

    result = PPTParser().parse(str(out))

    assert result["metadata"]["type"] == "ppt"
    assert result["metadata"]["slide_count"] == 3
    assert len(result["chunks"]) == 3
    for idx, chunk in enumerate(result["chunks"]):
        assert chunk["strategy"] == "ppt_slide"
        # page_number is 1-indexed per spec §3.2
        assert chunk["page_number"] == idx + 1
        assert chunk["modality"] == "text"


def test_ppt_parser_slide_chunk_text_contains_title_and_body(tmp_path):
    out = tmp_path / "named.pptx"
    prs = _new_pptx()
    _add_text_slide(prs, "Hello", "World")
    prs.save(str(out))

    result = PPTParser().parse(str(out))
    chunk_text = result["chunks"][0]["content"]
    # Format: "Slide 1\n<title>\n<body>"
    assert chunk_text.startswith("Slide 1")
    assert "Hello" in chunk_text
    assert "World" in chunk_text


def test_ppt_parser_speaker_notes_become_separate_chunk(tmp_path):
    """Notes share the slide's page_number so search can scope to a
    specific slide ("find me notes for slide 2")."""
    out = tmp_path / "notes.pptx"
    prs = _new_pptx()
    _add_text_slide(prs, "Slide 1", "Body", notes="Notes for slide 1")
    _add_text_slide(prs, "Slide 2", "Body", notes="Notes for slide 2")
    prs.save(str(out))

    result = PPTParser().parse(str(out))

    # 2 slides + 2 notes = 4 chunks
    assert len(result["chunks"]) == 4

    notes_chunks = [c for c in result["chunks"] if c["strategy"] == "ppt_notes"]
    assert len(notes_chunks) == 2

    # Notes for slide 1 has page_number=1 and contains its text
    slide1_notes = next(c for c in notes_chunks if c["page_number"] == 1)
    assert "Notes for slide 1" in slide1_notes["content"]
    assert slide1_notes["content"].startswith("Slide 1 Notes")

    slide2_notes = next(c for c in notes_chunks if c["page_number"] == 2)
    assert "Notes for slide 2" in slide2_notes["content"]


def test_ppt_parser_skips_notes_when_empty(tmp_path):
    """A slide with no notes_text_frame content produces NO notes chunk
    (zero-length text would poison the vector store)."""
    out = tmp_path / "no_notes.pptx"
    prs = _new_pptx()
    _add_text_slide(prs, "Slide", "Body")  # no notes
    prs.save(str(out))

    result = PPTParser().parse(str(out))
    # Only the slide chunk, no notes chunk.
    notes_chunks = [c for c in result["chunks"] if c["strategy"] == "ppt_notes"]
    assert notes_chunks == []


def test_ppt_parser_extracts_images_into_image_assets(tmp_path):
    """Images embedded in a slide become ``image_assets`` entries with
    mime + size + deterministic dedup key. The parser returns raw
    bytes; the caller is responsible for handing them to storage +
    inserting an ``ImageAsset`` row."""
    out = tmp_path / "with_img.pptx"
    prs = _new_pptx()
    _add_text_slide(prs, "Intro", "")
    _add_image_slide(prs, color="blue")
    _add_text_slide(prs, "Outro", "")
    prs.save(str(out))

    result = PPTParser().parse(str(out))

    # 3 slide chunks (Intro / image / Outro), no notes chunks, 1 image asset
    assert len(result["chunks"]) == 3
    assert "image_assets" in result
    assert len(result["image_assets"]) == 1

    asset = result["image_assets"][0]
    assert asset["page_number"] == 2  # the image slide is slide 2
    assert asset["mime"] == "image/png"
    assert asset["extension"] == "png"
    assert asset["size_bytes"] > 0
    assert isinstance(asset["bytes"], (bytes, bytearray))
    # Dedup key is stable per (slide, shape_index) so re-uploads of the
    # same PPT produce the same storage key.
    assert asset["slide_dedup_key"] == "slide_2_img_1"


def test_ppt_parser_image_count_in_metadata(tmp_path):
    out = tmp_path / "many_imgs.pptx"
    prs = _new_pptx()
    _add_image_slide(prs, "red")
    _add_image_slide(prs, "green")
    _add_image_slide(prs, "blue")
    prs.save(str(out))

    result = PPTParser().parse(str(out))
    assert result["metadata"]["image_count"] == 3


def test_ppt_parser_preserves_chunks_flag_is_true():
    assert PPTParser.preserves_chunks is True


def test_ppt_parser_get_type():
    assert PPTParser().get_type() == "ppt"


def test_ppt_parser_handles_missing_file(tmp_path):
    result = PPTParser().parse(str(tmp_path / "missing.pptx"))
    assert result["metadata"]["fallback"] is True
    assert "parse_error" in result["metadata"]


def test_ppt_parser_handles_corrupt_pptx(tmp_path):
    """A file with .pptx extension that isn't real pptx must fall back
    to the legacy path rather than raising."""
    out = tmp_path / "fake.pptx"
    out.write_bytes(b"this is not a pptx")
    result = PPTParser().parse(str(out))
    assert result["metadata"]["fallback"] is True
    assert "parse_error" in result["metadata"]


def test_ppt_parser_empty_presentation_produces_zero_chunks(tmp_path):
    """An empty pptx (no slides at all) yields zero chunks — there's
    nothing to embed. Downstream code handles zero-chunk case."""
    out = tmp_path / "empty.pptx"
    _new_pptx().save(str(out))

    result = PPTParser().parse(str(out))
    assert result["metadata"]["slide_count"] == 0
    assert result["chunks"] == []
    assert result["metadata"]["chunk_count"] == 0


def test_ppt_parser_page_numbers_are_one_indexed(tmp_path):
    """The spec (§3.2) is explicit: page_number=1 means the FIRST slide,
    not index 0. Verify by inspecting a 3-slide pptx."""
    out = tmp_path / "page_nums.pptx"
    prs = _new_pptx()
    for i in range(3):
        _add_text_slide(prs, f"Slide {i+1}", "")
    prs.save(str(out))

    result = PPTParser().parse(str(out))
    page_numbers = [c["page_number"] for c in result["chunks"]]
    assert page_numbers == [1, 2, 3]
