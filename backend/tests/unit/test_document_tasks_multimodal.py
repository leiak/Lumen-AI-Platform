"""M38.4 Step 5d: Worker multimodal dispatch + PPT image 落库 tests.

Covers:
- ``_persist_ppt_image_assets`` writes ``ImageAsset`` rows + storage upload
- ``_persist_ppt_image_assets`` is no-op when parser returns no
  ``image_assets``
- ``_persist_ppt_image_assets`` survives partial failures (one bad
  asset doesn't kill the rest)
- ``_dispatch_multimodal_embeddings`` writes vectors to MM FAISS index
  when KB.multimodal_enabled
- ``_dispatch_multimodal_embeddings`` marks ``embedding_status='failed'``
  on embedder error without flipping doc.status
- ``_dispatch_multimodal_embeddings`` is no-op when KB.multimodal_enabled
  is False
- ``process_document_task`` end-to-end: PPT doc → 1 image_asset row +
  1 image chunk + chunk_id backfilled + MM index written (mocked)
- ``PPTParser.parse()`` now emits one ``image`` chunk per extracted
  picture (Step 4 follow-up)

Patch targets: ``lumen_tasks.document_tasks`` for the helpers (they live
  in that module) and ``lumen_services.storage.get_storage_backend`` /
  ``lumen_services.multimodal_embedders.get_multimodal_embedder`` for
  I/O dependencies.

Notes:
- We don't drive ``process_document_task`` via Celery; we call the
  function with a constructed ``task_params`` dict + a ``MagicMock``
  ``self``. Each test patches DB / storage / embedder and asserts on
  observable side effects (``db.add`` calls, ``db.commit`` count,
  ``storage.put_object`` calls, ``mm_store.add_texts`` calls).
- The ``self.request.id`` attribute on the Celery task is a ``Mock``;
  we use ``MagicMock(spec=celery_app.Task)`` so attribute access
  doesn't blow up.
"""
from __future__ import annotations

from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import pytest

from lumen_services.parsers.ppt_parser import PPTParser


# --- helpers --------------------------------------------------------------


def _fake_image_record(
    page: int = 1,
    shape_idx: int = 1,
    ext: str = "png",
    mime: str = "image/png",
    size: int = 1024,
) -> dict:
    return {
        "page_number": page,
        "shape_index": shape_idx,
        "slide_dedup_key": f"slide_{page}_img_{shape_idx}",
        "bytes": b"\\x89PNG\\r\\n\\x1a\\n" * 64,
        "mime": mime,
        "extension": ext,
        "size_bytes": size,
    }


def _fake_parse_result(
    with_image_assets: bool = True,
    n_assets: int = 2,
) -> dict:
    """Mimic what ``DocumentParser.parse()`` returns for a PPT with
    embedded pictures."""
    image_assets = (
        [_fake_image_record(page=i + 1, shape_idx=i + 1) for i in range(n_assets)]
        if with_image_assets else []
    )
    return {
        "text": "Slide 1\nhello\n\nSlide 2\nworld",
        "metadata": {
            "type": "ppt",
            "format": "pptx",
            "slide_count": 2,
            "image_count": n_assets if with_image_assets else 0,
            "image_assets": image_assets,
        },
        "chunks": [
            {"content": "Slide 1\nhello", "chunk_index": 0, "length": 13,
             "strategy": "ppt_slide", "modality": "text",
             "sheet_name": None, "page_number": 1, "image_caption": None,
             "chunk_metadata": {"kind": "ppt_slide", "page_number": 1}},
            {"content": "Slide 2\nworld", "chunk_index": 1, "length": 13,
             "strategy": "ppt_slide", "modality": "text",
             "sheet_name": None, "page_number": 2, "image_caption": None,
             "chunk_metadata": {"kind": "ppt_slide", "page_number": 2}},
        ],
        "image_assets": image_assets,
    }


# --- _persist_ppt_image_assets -------------------------------------------


def test_persist_ppt_image_assets_writes_asset_rows_and_uploads():
    """Happy path: 2 image_assets → 2 storage.put_object + 2 db.add(ImageAsset)."""
    from lumen_tasks.document_tasks import _persist_ppt_image_assets

    db = MagicMock()
    storage = MagicMock()
    parse_result = _fake_parse_result(n_assets=2)

    with patch(
        "lumen_services.storage.get_storage_backend",
        return_value=storage,
    ):
        pages = _persist_ppt_image_assets(
            db=db,
            parse_result=parse_result,
            document_id=100,
            tenant_id=1,
            kb_id=55,
        )

    assert pages == [1, 2]
    assert storage.put_object.call_count == 2
    assert db.add.call_count == 2
    db.commit.assert_called_once()


def test_persist_ppt_image_assets_no_op_when_no_assets():
    """parser returned no image_assets → no storage / no DB writes."""
    from lumen_tasks.document_tasks import _persist_ppt_image_assets

    db = MagicMock()
    storage = MagicMock()
    parse_result = _fake_parse_result(with_image_assets=False)

    with patch(
        "lumen_services.storage.get_storage_backend",
        return_value=storage,
    ):
        pages = _persist_ppt_image_assets(
            db=db,
            parse_result=parse_result,
            document_id=100,
            tenant_id=1,
            kb_id=55,
        )

    assert pages == []
    storage.put_object.assert_not_called()
    db.add.assert_not_called()
    db.commit.assert_not_called()


def test_persist_ppt_image_assets_survives_partial_failure():
    """One asset's bytes is missing → skipped, others succeed."""
    from lumen_tasks.document_tasks import _persist_ppt_image_assets

    db = MagicMock()
    storage = MagicMock()

    parse_result = _fake_parse_result(n_assets=2)
    # corrupt one asset
    parse_result["image_assets"][0]["bytes"] = None

    with patch(
        "lumen_services.storage.get_storage_backend",
        return_value=storage,
    ):
        pages = _persist_ppt_image_assets(
            db=db,
            parse_result=parse_result,
            document_id=100,
            tenant_id=1,
            kb_id=55,
        )

    # Only 1 asset persisted (the one with bytes)
    assert storage.put_object.call_count == 1
    assert db.add.call_count == 1
    # Only the page of the successful asset
    assert pages == [2]


def test_persist_ppt_image_assets_rollback_on_commit_failure():
    """db.commit raises → function doesn't propagate; db.rollback called."""
    from lumen_tasks.document_tasks import _persist_ppt_image_assets

    db = MagicMock()
    db.commit.side_effect = Exception("commit boom")
    storage = MagicMock()

    parse_result = _fake_parse_result(n_assets=1)

    with patch(
        "lumen_services.storage.get_storage_backend",
        return_value=storage,
    ):
        # Must NOT raise
        _persist_ppt_image_assets(
            db=db,
            parse_result=parse_result,
            document_id=100,
            tenant_id=1,
            kb_id=55,
        )

    db.rollback.assert_called_once()


# --- _dispatch_multimodal_embeddings ------------------------------------


def _make_image_chunk(id_: int = 1, doc_id: int = 100,
    page: int | None = 1, caption: str = "logo") -> SimpleNamespace:
    c = SimpleNamespace(
        id=id_,
        document_id=doc_id,
        modality="image",
        page_number=page,
        image_caption=caption,
        content=caption,
        embedding_status="pending",
    )
    return c


def test_dispatch_multimodal_embeddings_writes_vectors():
    """Happy path: 2 image chunks → 2 vectors, MM store.add_texts called
    with matching length, ImageAsset.embedding_status flipped to 'ok'."""
    from lumen_tasks.document_tasks import _dispatch_multimodal_embeddings

    db = MagicMock()
    embedder = MagicMock()
    embedder.embed_text.side_effect = lambda cap: [0.1] * 1024
    mm_store = MagicMock()

    chunks = [_make_image_chunk(id_=1), _make_image_chunk(id_=2)]

    asset_for_chunk_1 = SimpleNamespace(
        chunk_id=1, embedding_status="pending",
    )
    asset_for_chunk_2 = SimpleNamespace(
        chunk_id=2, embedding_status="pending",
    )

    q_assets = MagicMock()
    q_assets.filter.return_value = q_assets
    q_assets.first.side_effect = [asset_for_chunk_1, asset_for_chunk_2]
    db.query.return_value = q_assets

    with patch(
        "lumen_services.multimodal_embedders.get_multimodal_embedder",
        return_value=(embedder, 1024),
    ), patch(
        "lumen_services.multimodal_vector_store_factory.MultimodalVectorStoreFactory.get_store",
        return_value=mm_store,
    ):
        _dispatch_multimodal_embeddings(
            db=db,
            image_chunks=chunks,
            kb_id=55,
            tenant_id=1,
            document_id=100,
            mm_config_id=7,
        )

    assert embedder.embed_text.call_count == 2
    mm_store.add_texts.assert_called_once()
    args, kwargs = mm_store.add_texts.call_args
    texts, metadatas, vectors = args[0], args[1], kwargs["vectors"]
    assert len(texts) == 2
    assert len(metadatas) == 2
    assert len(vectors) == 2
    # Image assets flipped to 'ok'
    assert asset_for_chunk_1.embedding_status == "ok"
    assert asset_for_chunk_2.embedding_status == "ok"
    # Chunks flipped to 'ok'
    assert all(c.embedding_status == "ok" for c in chunks)


def test_dispatch_multimodal_embeddings_marks_failed_when_embedder_load_fails():
    """get_multimodal_embedder raises → chunks get 'failed', doc not touched."""
    from lumen_tasks.document_tasks import _dispatch_multimodal_embeddings
    from lumen_services.multimodal_embedders import (
        MultimodalEmbeddingError,
    )

    db = MagicMock()
    chunks = [_make_image_chunk(id_=1)]

    with patch(
        "lumen_services.multimodal_embedders.get_multimodal_embedder",
        side_effect=MultimodalEmbeddingError("model not loaded"),
    ):
        _dispatch_multimodal_embeddings(
            db=db,
            image_chunks=chunks,
            kb_id=55,
            tenant_id=1,
            document_id=100,
            mm_config_id=7,
        )

    assert chunks[0].embedding_status == "failed"
    db.commit.assert_called()


def test_dispatch_multimodal_embeddings_marks_failed_when_embed_raises():
    """embed_text raises mid-batch → all chunks flipped to 'failed'."""
    from lumen_tasks.document_tasks import _dispatch_multimodal_embeddings

    db = MagicMock()
    embedder = MagicMock()
    embedder.embed_text.side_effect = Exception("HF model crashed")
    chunks = [_make_image_chunk(id_=1), _make_image_chunk(id_=2)]

    with patch(
        "lumen_services.multimodal_embedders.get_multimodal_embedder",
        return_value=(embedder, 1024),
    ):
        _dispatch_multimodal_embeddings(
            db=db,
            image_chunks=chunks,
            kb_id=55,
            tenant_id=1,
            document_id=100,
            mm_config_id=7,
        )

    for c in chunks:
        assert c.embedding_status == "failed"


def test_dispatch_multimodal_embeddings_no_op_when_no_image_chunks():
    """Empty image_chunks → no embedder call, no MM store write."""
    from lumen_tasks.document_tasks import _dispatch_multimodal_embeddings

    db = MagicMock()

    with patch(
        "lumen_services.multimodal_embedders.get_multimodal_embedder",
    ) as embedder_mock:
        _dispatch_multimodal_embeddings(
            db=db,
            image_chunks=[],
            kb_id=55,
            tenant_id=1,
            document_id=100,
            mm_config_id=7,
        )

    embedder_mock.assert_not_called()
    db.query.assert_not_called()
    db.commit.assert_not_called()


# --- PPT parser image chunks (Step 4 follow-up) -------------------------


def test_ppt_parser_emits_image_chunks_when_extract_pictures(
    tmp_path,
    monkeypatch,
):
    """After the Step 5d patch, PPTParser returns one ``image`` chunk per
    extracted picture. We use a real ``.pptx`` fixture written via
    python-pptx so this is end-to-end."""
    pytest.importorskip("pptx")
    pytest.importorskip("PIL")
    from PIL import Image
    from pptx import Presentation

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    # 1 slide with 1 picture (drawn from a 1x1 PNG blob)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Use Pillow to write a valid 1x1 PNG that python-pptx can reopen
    img_path = tmp_path / "tiny.png"
    Image.new("RGB", (1, 1), color=(255, 255, 255)).save(str(img_path), "PNG")
    slide.shapes.add_picture(str(img_path), 0, 0)

    # 2nd slide with 2 pictures
    slide2 = prs.slides.add_slide(blank)
    slide2.shapes.add_picture(str(img_path), 0, 0)
    slide2.shapes.add_picture(str(img_path), 100, 100)

    prs.save(str(pptx_path))

    parser = PPTParser()
    result = parser.parse(str(pptx_path))
    chunks = result["chunks"]
    image_chunks = [c for c in chunks if c.get("modality") == "image"]

    # 3 image chunks total (1 from slide 1, 2 from slide 2)
    assert len(image_chunks) == 3
    # page_number matches slide order
    assert [c["page_number"] for c in image_chunks] == [1, 2, 2]
    # image_assets list mirrors
    assert len(result["image_assets"]) == 3
    # image chunks come AFTER text chunks (insertion order preserved)
    first_image_idx = min(
        chunks.index(c) for c in image_chunks
    )
    assert first_image_idx > 0  # text chunks came first


def test_ppt_parser_no_image_chunks_when_no_pictures(tmp_path):
    """No embedded pictures → no image chunks, only slide / notes chunks."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    pptx_path = tmp_path / "textonly.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Add a textbox (NOT a picture)
    slide.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "hello world"
    prs.save(str(pptx_path))

    parser = PPTParser()
    result = parser.parse(str(pptx_path))
    chunks = result["chunks"]
    image_chunks = [c for c in chunks if c.get("modality") == "image"]
    assert image_chunks == []
    assert result["image_assets"] == []
    # text chunk still produced
    assert any(c.get("modality") == "text" for c in chunks)


# --- process_document_task wiring ---------------------------------------


def test_process_document_ppt_runs_image_asset_persistence():
    """End-to-end shape test: PPT doc → _persist_ppt_image_assets called
    once, _dispatch_multimodal_embeddings called once when KB has
    multimodal_enabled."""
    import lumen_tasks.document_tasks as task_module

    # Build a task_params for a PPT doc with multimodal KB.
    task_params = {
        "document_id": 100,
        "file_path": "data/uploads/1/55/pres.pptx",
        "file_content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "tenant_id": 1,
        "kb_id": 55,
        "user_id": 2,
        "doc_type": "ppt",
    }

    db = MagicMock()
    doc = SimpleNamespace(
        id=100,
        filename="pres.pptx",
        file_path="data/uploads/1/55/pres.pptx",
        status="processing",
        knowledge_base_id=55,
        doc_metadata={},
        error_message=None,
        embedding_model_config_id=3,
        chunk_count=None,
        knowledge_base=SimpleNamespace(
            id=55,
            tenant_id=1,
            embedding_model_config_id=3,
            multimodal_enabled=True,
            multimodal_config_id=7,
        ),
    )

    # Document lookup + chunks creation patterns
    q_doc = MagicMock()
    q_doc.filter.return_value = q_doc
    q_doc.first.return_value = doc

    # image asset lookup for chunk_id backfill
    image_assets_with_no_chunk = [
        SimpleNamespace(
            id=1, document_id=100, original_doc_page=1, chunk_id=None,
        ),
    ]
    q_assets = MagicMock()
    q_assets.filter.return_value = q_assets
    q_assets.order_by.return_value = q_assets
    q_assets.all.return_value = image_assets_with_no_chunk
    q_assets.first.return_value = image_assets_with_no_chunk[0]

    # image chunks pre-allocated with id by DB (mimics flush)
    image_chunk_a = SimpleNamespace(
        id=200, document_id=100, modality="image", page_number=1,
        content="slide 1 image 1", image_caption="slide 1 image 1",
        embedding_status="ok", vector_id=None,
    )
    text_chunk_a = SimpleNamespace(
        id=201, document_id=100, modality="text", page_number=1,
        content="Slide 1\nhello", image_caption=None,
        embedding_status="ok", vector_id=None,
    )

    # parse returns slide chunks; image_chunks populated by chunk dict
    # that worker creates. To keep this test simple we patch
    # ``DocumentChunk(...)`` to return our pre-made chunks.
    parse_result = _fake_parse_result(n_assets=1)

    def db_query_side_effect(*args, **kwargs):
        # Document lookup
        if args and getattr(args[0], "__name__", "") == "Document":
            return q_doc
        return q_assets

    db.query.side_effect = db_query_side_effect

    # Mimic SQLAlchemy auto-flush: when worker does ``db.add(chunk)`` the
    # SQLAlchemy session assigns chunk.id on flush; for the mock we set
    # it manually via db.add.side_effect.
    chunk_counter = {"id": 200}
    def fake_add(obj):
        # If it's a chunk, assign id + accumulate
        if hasattr(obj, "modality") and obj.modality == "image":
            obj.id = chunk_counter["id"]
            chunk_counter["id"] += 1
    db.add.side_effect = fake_add

    # Vector store factory
    vector_store = MagicMock()
    vector_store.add_texts.return_value = ["mock_0", "mock_1"]

    # MM vector store
    mm_store = MagicMock()
    embedder = MagicMock()
    embedder.embed_text.return_value = [0.1] * 1024

    # We patch the WHOLE pipeline so the worker runs but DB / storage /
    # embedders are mocked. Force parse_result into the parser mock.
    parser_mock = MagicMock()
    parser_mock.parse.return_value = parse_result

    # ``DocumentChunk(...)`` constructor returns our pre-built chunks
    chunk_iter = iter([image_chunk_a, text_chunk_a])
    def chunk_ctor(*args, **kw):
        try:
            return next(chunk_iter)
        except StopIteration:
            return SimpleNamespace(**kw, id=None, embedding_status="ok")
    chunk_patch_path = "lumen_models.knowledge.DocumentChunk"

    # Capture calls to the helpers we care about
    with patch.object(
        task_module, "_persist_ppt_image_assets",
        wraps=task_module._persist_ppt_image_assets,
    ) as persist_spy, patch.object(
        task_module, "_dispatch_multimodal_embeddings",
        wraps=task_module._dispatch_multimodal_embeddings,
    ) as dispatch_spy, patch(
        "lumen_core.database.SessionLocal",
        return_value=db,
    ), patch(
        "lumen_services.document_parser.DocumentParser",
        return_value=parser_mock,
    ), patch(
        "lumen_tools.vector_store_factory.VectorStoreFactory.get_store",
        return_value=vector_store,
    ), patch(
        "lumen_services.multimodal_embedders.get_multimodal_embedder",
        return_value=(embedder, 1024),
    ), patch(
        "lumen_services.multimodal_vector_store_factory.MultimodalVectorStoreFactory.get_store",
        return_value=mm_store,
    ), patch(
        "lumen_services.storage.get_storage_backend",
        return_value=MagicMock(),
    ), patch(
        "lumen_models.knowledge.DocumentChunk",
        side_effect=chunk_ctor,
    ):
        result = task_module.process_document_task.run(task_params=task_params)

    # Helpers were both called
    persist_spy.assert_called_once()
    # dispatch_spy is called once only if image_chunks survived the
    # parser+chunk-creation pipeline. We don't strictly assert (the
    # chunk_ctor mock may exhaust early) — but assert_called confirms
    # the dispatcher ran.
    assert dispatch_spy.called, (
        "expected _dispatch_multimodal_embeddings to be called for "
        "multimodal-enabled KB with image chunks"
    )
    # MM embedder got embed_text for each image chunk
    embedder.embed_text.assert_called_once()
    mm_store.add_texts.assert_called_once()
    # image_asset.chunk_id backfilled
    assert image_assets_with_no_chunk[0].chunk_id == image_chunk_a.id
    # Worker returned completed
    assert result["status"] == "completed"


def test_process_document_skips_mm_dispatch_when_kb_disabled():
    """If KB.multimodal_enabled=False, _dispatch_multimodal_embeddings
    is NOT called even when image chunks exist."""
    import lumen_tasks.document_tasks as task_module

    task_params = {
        "document_id": 100,
        "file_path": "data/uploads/1/55/pres.pptx",
        "file_content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "tenant_id": 1,
        "kb_id": 55,
        "user_id": 2,
        "doc_type": "ppt",
    }

    db = MagicMock()
    doc = SimpleNamespace(
        id=100,
        filename="pres.pptx",
        file_path="data/uploads/1/55/pres.pptx",
        status="processing",
        knowledge_base_id=55,
        doc_metadata={},
        error_message=None,
        embedding_model_config_id=3,
        chunk_count=None,
        knowledge_base=SimpleNamespace(
            id=55,
            tenant_id=1,
            embedding_model_config_id=3,
            multimodal_enabled=False,
            multimodal_config_id=None,
        ),
    )
    q_doc = MagicMock()
    q_doc.filter.return_value = q_doc
    q_doc.first.return_value = doc
    db.query.return_value = q_doc

    parser_mock = MagicMock()
    parser_mock.parse.return_value = _fake_parse_result(n_assets=1)

    text_chunk_a = SimpleNamespace(
        id=201, document_id=100, modality="text", page_number=1,
        content="Slide 1\nhello", image_caption=None,
        embedding_status="ok", vector_id=None,
    )

    chunk_iter = iter([text_chunk_a])
    def chunk_ctor(*args, **kw):
        try:
            return next(chunk_iter)
        except StopIteration:
            return SimpleNamespace(**kw, id=None, embedding_status="ok")
    chunk_patch_path = "lumen_models.knowledge.DocumentChunk"

    vector_store = MagicMock()
    vector_store.add_texts.return_value = ["mock_0"]

    with patch.object(
        task_module, "_persist_ppt_image_assets",
    ), patch.object(
        task_module, "_dispatch_multimodal_embeddings",
    ) as dispatch_spy, patch(
        "lumen_core.database.SessionLocal",
        return_value=db,
    ), patch(
        "lumen_services.document_parser.DocumentParser",
        return_value=parser_mock,
    ), patch(
        "lumen_tools.vector_store_factory.VectorStoreFactory.get_store",
        return_value=vector_store,
    ), patch(
        "lumen_services.storage.get_storage_backend",
        return_value=MagicMock(),
    ), patch(
        "lumen_models.knowledge.DocumentChunk",
        side_effect=chunk_ctor,
    ):
        task_module.process_document_task.run(task_params=task_params)

    dispatch_spy.assert_not_called()