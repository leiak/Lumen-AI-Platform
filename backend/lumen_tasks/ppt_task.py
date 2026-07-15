"""Celery task for PPT generation (python-pptx high-precision mode).

Spec: docs-internal/superpowers/specs/m35-ppt-generation.md §8
"""
import logging
import uuid
from dataclasses import dataclass
from pathlib import Path

from lumen_tasks.celery_app import celery_app
from lumen_core.config import settings
from lumen_core.database import SessionLocal
from lumen_models.ppt_task import PptTask
from lumen_schemas.ppt import PptSchema
from lumen_services.ppt_service import PptService
from lumen_services.notification_service import NotificationService
from lumen_services.electron_service import broadcast_event_sync

logger = logging.getLogger(__name__)

# Pre-load SQLAlchemy models to avoid Windows spawn multiprocessing issue:
# when Celery pickles+spawns a fresh process, module-level relationship()
# string references (e.g. KnowledgeBase.tenant) fail if Tenant wasn't yet
# registered with the SQLAlchemy registry.  Import in dependency order.
import lumen_models.tenant  # noqa: F401
import lumen_models.user  # noqa: F401
import lumen_models.knowledge  # noqa: F401
import lumen_models.chat  # noqa: F401
import lumen_models.agent  # noqa: F401
import lumen_models.agent_team  # noqa: F401  # AgentTeam (LLMCallLog.team_id FK)
import lumen_models.embedding_call_log  # noqa: F401  # EmbeddingCallLog
import lumen_models.llm_call_log  # noqa: F401  # LLMCallLog (WorkflowRun relationship)
import lumen_models.workflow  # noqa: F401  # WorkflowRun references EmbeddingCallLog
import lumen_models.model_config  # noqa: F401  # ModelConfig (LLMCallLog.model_config_id FK)
import lumen_models.image_generation  # noqa: F401  # GeneratedImage (LLMCallLog.image_id FK)


@dataclass(frozen=True)
class StyleTheme:
    """一个风格的完整视觉主题（颜色 / 字体 / 装饰 / 图表色板 / 页面 chrome）。"""
    bg: tuple            # 背景 RGB
    title: tuple         # 标题 RGB
    body: tuple          # 正文 RGB
    accent: tuple        # 装饰条 / 竖条 RGB
    accent_soft: tuple   # 浅装饰 RGB
    title_font: str      # 标题字体
    body_font: str       # 正文字体
    title_size: int      # 内容页标题字号 pt
    body_size: int       # 正文字号 pt
    bullet_char: str     # 项目符号字符
    chart_palette: tuple # 图表色板（hex 字符串）
    show_top_bar: bool   # 顶部 / 底部装饰条
    show_left_accent: bool  # 内容页左侧竖条 + 框
    show_decorative_line: bool  # 封面标题下方装饰线
    page_number_format: str  # "dotted"(01.) / "slash"(3/8) / "roman"(III / IX)


_THEMES = {
    "simple": StyleTheme(
        bg=(0xFF, 0xFF, 0xFF),
        title=(0x1A, 0x1A, 0x1A),
        body=(0x44, 0x44, 0x44),
        accent=(0x88, 0x88, 0x88),
        accent_soft=(0xCC, 0xCC, 0xCC),
        title_font="Calibri Light",
        body_font="Calibri Light",
        title_size=30,
        body_size=20,
        bullet_char="—",
        chart_palette=("595959", "808080", "A6A6A6", "BFBFBF", "D9D9D9", "8C8C8C"),
        show_top_bar=False,
        show_left_accent=False,
        show_decorative_line=False,
        page_number_format="dotted",
    ),
    "business": StyleTheme(
        bg=(0xFF, 0xFF, 0xFF),
        title=(0x1F, 0x4E, 0x79),
        body=(0x33, 0x33, 0x33),
        accent=(0x2C, 0x5A, 0x8A),
        accent_soft=(0x44, 0x72, 0xC4),
        title_font="Calibri",
        body_font="Calibri",
        title_size=30,
        body_size=20,
        bullet_char="▪",
        chart_palette=("1F4E79", "ED7D31", "70AD47", "FFC000", "4472C4", "264478"),
        show_top_bar=True,
        show_left_accent=True,
        show_decorative_line=True,
        page_number_format="slash",
    ),
    "academic": StyleTheme(
        bg=(0xFA, 0xFA, 0xF5),
        title=(0x00, 0x32, 0x8B),
        body=(0x22, 0x22, 0x22),
        accent=(0x00, 0x32, 0x8B),
        accent_soft=(0x80, 0x80, 0x00),
        title_font="Times New Roman",
        body_font="Times New Roman",
        title_size=30,
        body_size=20,
        bullet_char="§",
        chart_palette=("4682B4", "B22222", "808000", "CD853F", "2F4F4F", "8B4513"),
        show_top_bar=True,
        show_left_accent=False,
        show_decorative_line=False,
        page_number_format="roman",
    ),
}


def _style_theme(style: str) -> StyleTheme:
    """返回风格对应的 StyleTheme（未知风格回落到 simple）。"""
    return _THEMES.get(style, _THEMES["simple"])


def _to_roman(n: int) -> str:
    vals = [
        (1000, "M"), (900, "CM"), (500, "D"), (400, "CD"), (100, "C"),
        (90, "XC"), (50, "L"), (40, "XL"), (10, "X"), (9, "IX"),
        (5, "V"), (4, "IV"), (1, "I"),
    ]
    out = []
    for v, s in vals:
        while n >= v:
            out.append(s)
            n -= v
    return "".join(out) or "I"


def _format_page_number(fmt: str, num: int, total: int) -> str:
    if fmt == "dotted":
        return f"{num:02d}."
    if fmt == "roman":
        return f"{_to_roman(num)} / {_to_roman(total)}"
    return f"{num}/{total}"  # slash（默认）


def _render_pptx(schema: PptSchema, output_path: Path, style: str) -> None:
    """用 python-pptx 渲染 PPT 文件，按 style 主题差异化视觉。"""
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData

    theme = _style_theme(style)
    bg_color = RGBColor(*theme.bg)
    title_color = RGBColor(*theme.title)
    body_color = RGBColor(*theme.body)
    accent_color = RGBColor(*theme.accent)
    accent_soft_color = RGBColor(*theme.accent_soft)
    content_bg_color = RGBColor(0xF5, 0xF5, 0xF5)
    bullet = theme.bullet_char

    def _apply_font(paragraph, font_name: str):
        try:
            paragraph.font.name = font_name
        except Exception:
            pass

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    total = len(schema.slides)

    for idx, slide_data in enumerate(schema.slides):
        slide = prs.slides.add_slide(blank_layout)

        # 背景色
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = bg_color

        slide_num = idx + 1

        # ===== 顶部 / 底部装饰条（依主题）=====
        if theme.show_top_bar:
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = accent_color
            top_bar.line.fill.background()

            # academic：顶部双线（第二条 0.05" 细线）
            if theme.page_number_format == "roman":
                top_bar2 = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.16), prs.slide_width, Inches(0.05)
                )
                top_bar2.fill.solid()
                top_bar2.fill.fore_color.rgb = accent_color
                top_bar2.line.fill.background()

            bottom_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), prs.slide_width, Inches(0.12)
            )
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = accent_color
            bottom_bar.line.fill.background()

        # ===== 页码（右下角）=====
        page_box = slide.shapes.add_textbox(
            Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.35)
        )
        page_tf = page_box.text_frame
        page_p = page_tf.paragraphs[0]
        page_p.text = _format_page_number(theme.page_number_format, slide_num, total)
        page_p.font.size = Pt(10)
        page_p.font.color.rgb = body_color
        page_p.alignment = PP_ALIGN.RIGHT
        _apply_font(page_p, theme.body_font)

        # ===== 封面页 =====
        if slide_data.layout == "title_only":
            main_title = slide_data.title or schema.title
            # 主标题（大号）
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.2), Inches(11.333), Inches(1.4)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = main_title
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.CENTER
            _apply_font(p, theme.title_font)

            # 标题下方装饰线（依主题）
            if theme.show_decorative_line:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.65), Inches(2.333), Inches(0.06)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = accent_soft_color
                line.line.fill.background()

            y_offset = 3.9
            if schema.subtitle:
                sub_box = slide.shapes.add_textbox(
                    Inches(1), Inches(3.9), Inches(11.333), Inches(0.7)
                )
                tf2 = sub_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = schema.subtitle
                p2.font.size = Pt(24)
                p2.font.color.rgb = body_color
                p2.alignment = PP_ALIGN.CENTER
                _apply_font(p2, theme.body_font)
                y_offset = 4.8

            # 作者信息
            author_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.8), Inches(11.333), Inches(0.5)
            )
            tf3 = author_box.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = schema.author or "Lumen AI"
            p3.font.size = Pt(16)
            p3.font.color.rgb = body_color
            p3.alignment = PP_ALIGN.CENTER
            _apply_font(p3, theme.body_font)

            # content（如果有）
            if slide_data.content:
                content_box = slide.shapes.add_textbox(
                    Inches(1), Inches(y_offset + 0.3), Inches(11.333), Inches(1.5)
                )
                tf4 = content_box.text_frame
                tf4.word_wrap = True
                for i, item in enumerate(slide_data.content):
                    p4 = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
                    p4.text = f"{bullet} {item}"
                    p4.font.size = Pt(16)
                    p4.font.color.rgb = body_color
                    p4.alignment = PP_ALIGN.CENTER
                    _apply_font(p4, theme.body_font)

        # ===== 结束页 =====
        elif slide_data.layout == "blank":
            end_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.8), Inches(11.333), Inches(1.5)
            )
            tf = end_box.text_frame
            p = tf.paragraphs[0]
            p.text = "谢谢观看"
            p.font.size = Pt(56)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.CENTER
            _apply_font(p, theme.title_font)

            # 装饰线（依主题）
            if theme.show_decorative_line:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.35), Inches(3.333), Inches(0.06)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = accent_soft_color
                line.line.fill.background()

            if schema.author:
                author_box = slide.shapes.add_textbox(
                    Inches(1), Inches(4.6), Inches(11.333), Inches(0.5)
                )
                tf2 = author_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = schema.author
                p2.font.size = Pt(16)
                p2.font.color.rgb = body_color
                p2.alignment = PP_ALIGN.CENTER
                _apply_font(p2, theme.body_font)

        # ===== 图表页 =====
        elif slide_data.layout == "chart":
            if slide_data.title:
                t_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
                )
                t_tf = t_box.text_frame
                t_p = t_tf.paragraphs[0]
                t_p.text = slide_data.title
                t_p.font.size = Pt(32)
                t_p.font.bold = True
                t_p.font.color.rgb = title_color
                _apply_font(t_p, theme.title_font)

            if slide_data.chart:
                chart_data = CategoryChartData()
                chart_data.categories = slide_data.chart.labels
                for ds in slide_data.chart.datasets:
                    chart_data.add_series(ds["name"], ds["values"])

                chart_type_map = {
                    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE,
                    "pie": XL_CHART_TYPE.PIE,
                }
                chart_type = chart_type_map.get(slide_data.chart.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
                chart_frame = slide.shapes.add_chart(
                    chart_type, Inches(1), Inches(1.4), Inches(11), Inches(5.5), chart_data
                )
                chart = chart_frame.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM

                # 主题图表色板
                try:
                    plot = chart.plots[0]
                    if slide_data.chart.type == "pie":
                        points = plot.series[0].points
                        for i, point in enumerate(points):
                            point.format.fill.solid()
                            point.format.fill.fore_color.rgb = RGBColor.from_string(
                                theme.chart_palette[i % len(theme.chart_palette)]
                            )
                    else:
                        for i, series in enumerate(plot.series):
                            series.format.fill.solid()
                            series.format.fill.fore_color.rgb = RGBColor.from_string(
                                theme.chart_palette[i % len(theme.chart_palette)]
                            )
                except Exception:
                    pass

        # ===== 内容页（title_content / two_column）=====
        else:
            # 左侧强调竖条（依主题）
            if theme.show_left_accent:
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(0.08), Inches(1.0)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = accent_color
                accent_bar.line.fill.background()

            # 标题
            if slide_data.title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.9)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.title
                p.font.size = Pt(theme.title_size)
                p.font.bold = True
                p.font.color.rgb = title_color
                _apply_font(p, theme.title_font)

            # 内容区域背景框（依主题：仅 business 加框）
            if theme.show_left_accent:
                content_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.4), Inches(1.4), Inches(12.533), Inches(5.7)
                )
                content_shape.fill.solid()
                content_shape.fill.fore_color.rgb = content_bg_color
                content_shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                content_shape.line.width = Pt(1)

            if slide_data.layout == "title_content" and slide_data.content:
                content_box = slide.shapes.add_textbox(
                    Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.3)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                for i, item in enumerate(slide_data.content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"{bullet} {item}"
                    p.font.size = Pt(theme.body_size)
                    p.font.color.rgb = body_color
                    p.space_after = Pt(14)
                    _apply_font(p, theme.body_font)

            elif slide_data.layout == "two_column":
                left = slide_data.leftContent or []
                right = slide_data.rightContent or []

                # 左栏
                left_box = slide.shapes.add_textbox(
                    Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.3)
                )
                tf_l = left_box.text_frame
                tf_l.word_wrap = True
                for i, item in enumerate(left):
                    p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
                    p.text = f"{bullet} {item}"
                    p.font.size = Pt(theme.body_size - 2)
                    p.font.color.rgb = body_color
                    p.space_after = Pt(10)
                    _apply_font(p, theme.body_font)

                # 右栏
                right_box = slide.shapes.add_textbox(
                    Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.3)
                )
                tf_r = right_box.text_frame
                tf_r.word_wrap = True
                for i, item in enumerate(right):
                    p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
                    p.text = f"{bullet} {item}"
                    p.font.size = Pt(theme.body_size - 2)
                    p.font.color.rgb = body_color
                    p.space_after = Pt(10)
                    _apply_font(p, theme.body_font)

        # 演讲者备注
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes

    prs.save(str(output_path))


@celery_app.task(bind=True, name="lumen_tasks.ppt_task.generate_ppt_task")
def generate_ppt_task(
    self,
    task_id: str,
    tenant_id: int,
    user_id: int,
    conversation_id: int,
    title: str,
    content_range: int,
    include_charts: bool,
    style: str,
):
    """Celery task: LLM 生成 PPT Schema → python-pptx 渲染 → 存文件。"""
    db = SessionLocal()
    try:
        # 更新状态为 processing
        task = db.query(PptTask).filter(PptTask.task_id == task_id).first()
        if not task:
            logger.warning("PptTask %s not found", task_id)
            return
        task.status = "processing"
        db.commit()

        # 1. 调用 LLM 生成 PPT Schema
        service = PptService()
        schema = service.build_schema(
            db=db,
            tenant_id=tenant_id,
            user_id=user_id,
            conversation_id=conversation_id,
            title=title,
            content_range=content_range,
            include_charts=include_charts,
            style=style,
        )

        # 2. 保存 schema_json（便于排查）
        import json as json_mod
        task.schema_json = json_mod.dumps(schema.model_dump(), ensure_ascii=False)
        db.commit()

        # 3. python-pptx 渲染
        output_dir = settings.STORAGE_DIR / "ppt" / str(tenant_id)
        output_dir.mkdir(parents=True, exist_ok=True)
        filename = f"{uuid.uuid4().hex}.pptx"
        output_path = output_dir / filename

        _render_pptx(schema, output_path, style)

        # 4. 更新任务状态
        file_url = f"/ppt/{tenant_id}/{filename}"
        task.status = "completed"
        task.file_url = file_url
        db.commit()

        # 5. 发通知
        try:
            NotificationService.publish_event(
                db,
                user_id=user_id,
                type="ppt_generation_completed",
                title="PPT 已生成",
                body=f"「{title}」已生成完毕，点击下载",
                resource_type="ppt_task",
                resource_id=None,
                metadata={"file_url": file_url, "task_id": task_id},
            )
            # WebSocket 广播
            broadcast_event_sync(
                event="ppt_task_completed",
                payload={
                    "task_id": task_id,
                    "file_url": file_url,
                    "title": title,
                },
                target_user_id=user_id,
            )
        except Exception as notify_err:
            logger.warning("PPT notification broadcast failed: %s", notify_err)

        logger.info("PptTask %s completed: %s", task_id, file_url)

    except Exception as e:
        logger.error("PptTask %s failed: %s", task_id, e)
        task = db.query(PptTask).filter(PptTask.task_id == task_id).first()
        if task:
            task.status = "failed"
            task.error = str(e)
            db.commit()
        raise
    finally:
        db.close()
